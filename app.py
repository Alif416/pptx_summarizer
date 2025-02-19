from flask import Flask, request, render_template, jsonify
from pptx import Presentation
from groq import Groq

app = Flask(__name__)

# Initialize Groq client
client = Groq(api_key="GROQ_API_KEY")  
#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
def extract_text_from_pptx(file):
    """
    Extract text from a PPTX file.
    """
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def summarize_text(text, model_name="llama-3.3-70b-versatile"):
    """
    Summarize text using Groq's API.
    """
    response = client.chat.completions.create(
        messages=[{
            "role": "user",
            "content": f"""first analyze the full text . after that
            summarize it in points and focus on the main points. Also use simple language. start with the summary "The pptx file contains the following information"
            
            :\n\n{text}"""
        }],
        model=model_name, 
        temperature=0.3
    )
    print(text)
    return response.choices[0].message.content

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        # Check if a file was uploaded
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        
        if file and file.filename.endswith(".pptx"):
            # Extract text from PPTX
            text = extract_text_from_pptx(file)
            
            # Summarize text
            summary = summarize_text(text)
            
            # Return the summary
            return jsonify({"summary": summary})
    
    return render_template("index.html")

if __name__ == "__main__":
    app.run(debug=True)